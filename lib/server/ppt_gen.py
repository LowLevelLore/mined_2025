from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE

ppt_content = {'title': 'SUFIA: Language-Guided Augmented Dexterity for Robotic Surgical Assistants',
 'authors': ['Masoud Moghani',
  'Lars Doorenbos',
  'William Chung-Ho Panitch',
  'Sean Huver',
  'Mahdi Azizian',
  'Ken Goldberg',
  'Animesh Garg'],
 'institution': ['University of Toronto',
  'University of Bern',
  'University of California, Berkeley',
  'NVIDIA',
  'Georgia Institute of Technology'],
 'slides': [{'title': 'Title Slide',
   'bullet_points': ['SUFIA: Language-Guided Augmented Dexterity for Robotic Surgical Assistants',
    'Masoud Moghani, Lars Doorenbos, William Chung-Ho Panitch, Sean Huver, Mahdi Azizian, Ken Goldberg, Animesh Garg',
    'University of Toronto, University of Bern, UC Berkeley, NVIDIA, Georgia Tech'],
   'notes': 'Start with a visually appealing title slide.  Consider including a relevant image, such as a robotic arm performing a surgical task.'},
  {'title': 'Introduction',
   'bullet_points': ['Problem: Current robotic surgical assistants (RSAs) often require tedious teleoperation, leading to surgeon fatigue and limiting autonomy.',
    'Objectives: Develop SUFIA, a framework for natural language-guided augmented dexterity in RSAs.',
    'Motivation:  Enable learning-free, generalizable surgical sub-task execution with a human-in-the-loop safety mechanism. Overcome limitations of learning-based approaches which are computationally expensive, require extensive data, and lack generalizability.',
    'Improve human-robot interaction in surgery through natural language commands.'],
   'notes': 'Highlight the gap in current technology and how SUFIA addresses it.  Emphasize the benefits of a language-guided approach.'},
  {'title': 'Methods',
   'bullet_points': ['Methodology:  Uses a Large Language Model (LLM) for high-level planning and low-level control code generation.',
    'Datasets:  Simulated surgical environments (ORBIT-Surgical) and a physical da Vinci Research Kit (dVRK) platform.',
    'Experimental Setup:  Four simulated surgical sub-tasks and two physical sub-tasks (Needle Lift, Needle Handover).  Uses RGB-D camera for perception, with a trained segmentation network for physical experiments and instance segmentation in simulation.',
    'LLM: Primarily uses GPT-4 Turbo, with comparisons to other LLMs.'],
   'notes': 'Include a diagram illustrating the system architecture.  Mention the specific LLMs used and any pre-training details.'},
  {'title': 'Results',
   'bullet_points': ['Success rates for simulated tasks (Table I): Needle Lift (100%), Needle Handover (90%), Vessel Dilation (60%), Shunt Insertion (70%).',
    'Success rates for physical tasks (Table I): Needle Lift (100%), Needle Handover (50%).',
    'Table I shows success rates and failure modes (planning vs. execution).',
    'Figure 3 shows the four simulated surgical sub-tasks.',
    'Figure 4 shows the physical Needle Handover task.',
    'Figure 5 shows variations in simulated needles.',
    'Table II shows the robustness of the perception module to different needle shapes and sizes in simulation.',
    'Analysis of task prompts:  Simple prompts for simple tasks, more complex prompts for multi-step tasks.  Examples of prompts provided.'],
   'images': ['table_I.png',
    'figure_3.png',
    'figure_4.png',
    'figure_5.png',
    'table_II.png'],
   'notes': 'Present results clearly and concisely.  Use charts and graphs to visualize key findings.  Discuss any limitations or unexpected results.'},
  {'title': 'Discussion',
   'bullet_points': ['Significance: SUFIA demonstrates a novel, learning-free approach to surgical augmented dexterity, generalizing across multiple tasks.',
    'Comparison with prior work:  Most prior work relies on task-specific models or extensive training data. SUFIA leverages the general reasoning capabilities of LLMs.',
    'Robustness:  SUFIA shows robustness to variations in object shape, size, and workspace conditions (both simulated and physical).',
    'Limitations:  Current real-time performance is limited by the LLM API call speed.  Potential risks remain from unexpected circumstances.'],
   'notes': "Compare SUFIA's performance and capabilities to existing methods. Discuss the implications of the findings."},
  {'title': 'Conclusion',
   'bullet_points': ['SUFIA is a modular framework for natural surgeon-robot interaction, successfully automating surgical sub-tasks using a learning-free approach.',
    'Safety is enhanced through re-planning and a human-in-the-loop design.',
    'Future work: Explore on-device execution of quantized open-source LLMs to improve speed and address privacy concerns. Investigate fine-tuned large language and vision models.'],
   'notes': 'Summarize the key contributions and potential impact of the research.  Outline future research directions.'},
  {'title': 'References',
   'bullet_points': ['[1] ...', '[2] ...', '...[47]'],
   'notes': 'Include a complete list of cited references in a consistent format.'}]}

class ThemeConfig:
    def __init__(self, name="modern"):
        themes = {
            "modern": {
                "background": RGBColor(30, 30, 30),  # Dark gray background
                "title": RGBColor(255, 215, 0),  # Gold title text
                "body": RGBColor(200, 200, 200),  # Light gray body text
                "title_font": "Montserrat",
                "body_font": "Lato",
            },
            "vintage": {
                "background": RGBColor(245, 222, 179),  # Wheat background
                "title": RGBColor(139, 69, 19),  # Saddle brown title
                "body": RGBColor(105, 105, 105),  # Dim gray text
                "title_font": "Georgia",
                "body_font": "Times New Roman",
            },
            "corporate": {
                "background": RGBColor(255, 255, 255),  # White background
                "title": RGBColor(0, 51, 102),  # Navy blue title
                "body": RGBColor(51, 51, 51),  # Dark gray body text
                "title_font": "Arial",
                "body_font": "Verdana",
            },
            "minimal": {
                "background": RGBColor(240, 240, 240),  # Light gray background
                "title": RGBColor(50, 50, 50),  # Dark gray title
                "body": RGBColor(80, 80, 80),  # Slightly lighter gray for body
                "title_font": "Helvetica",
                "body_font": "Sans-Serif",
            },
            "bold": {
                "background": RGBColor(0, 0, 0),  # Black background
                "title": RGBColor(255, 0, 0),  # Red title text
                "body": RGBColor(255, 255, 255),  # White body text
                "title_font": "Impact",
                "body_font": "Arial Black",
            }
        }
        self.theme = themes.get(name, themes["modern"])

def apply_background(slide, color):
    """Apply background color to a slide"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def create_themed_presentation(content, theme_name="default", output_file="presentation.pptx"):
    prs = Presentation()
    theme = ThemeConfig(theme_name)
    
    # Create title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    apply_background(title_slide, theme.theme["background"])
    
    # Set title
    title = title_slide.shapes.title
    title.text = content['title']
    title_para = title.text_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.name = theme.theme["title_font"]
    title_para.font.color.rgb = theme.theme["title"]
    title_para.alignment = PP_ALIGN.CENTER
    
    # Set subtitle (authors and institutions)
    subtitle = title_slide.placeholders[1]
    subtitle.text = (", ".join(content['authors']) + "\n" + 
                    ", ".join(content['institution']))
    sub_para = subtitle.text_frame.paragraphs[0]
    sub_para.font.size = Pt(18)
    sub_para.font.name = theme.theme["body_font"]
    sub_para.font.color.rgb = theme.theme["body"]
    sub_para.alignment = PP_ALIGN.CENTER
    
    # Create content slides
    for slide_data in content['slides'][1:]:  # Skip first slide (title)
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        apply_background(slide, theme.theme["background"])
        
        # Set slide title
        title = slide.shapes.title
        title.text = slide_data['title']
        title_para = title.text_frame.paragraphs[0]
        title_para.font.size = Pt(32)
        title_para.font.name = theme.theme["title_font"]
        title_para.font.color.rgb = theme.theme["title"]
        
        # Add bullet points
        content_placeholder = slide.placeholders[1]
        text_frame = content_placeholder.text_frame
        text_frame.clear()
        
        for point in slide_data['bullet_points']:
            paragraph = text_frame.add_paragraph()
            paragraph.text = point
            paragraph.font.size = Pt(18)
            paragraph.font.name = theme.theme["body_font"]
            paragraph.font.color.rgb = theme.theme["body"]
            
        # Add notes if present
        if 'notes' in slide_data:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = slide_data['notes']
            
    prs.save(output_file)
    print(f"Created presentation with {theme_name} theme: {output_file}")

# Example usage:
if __name__ == "__main__":
    create_themed_presentation(ppt_content, "modern", "sufia_modern_4.pptx")
