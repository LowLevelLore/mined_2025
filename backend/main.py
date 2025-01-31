from flask import Flask, request, send_file, redirect, jsonify
from flask_sqlalchemy import SQLAlchemy
import os
import io
import random
import pyttsx3
import glob
import wave
import zipfile
from flask_cors import CORS
from werkzeug.utils import secure_filename
import requests
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor


class ThemeConfig:
    def __init__(self, name="modern"):
        themes = {
            "modern": {
                "background": RGBColor(30, 30, 30),
                "title": RGBColor(255, 215, 0),
                "body": RGBColor(200, 200, 200),
                "title_font": "Montserrat",
                "body_font": "Lato",
            },
            "vintage": {
                "background": RGBColor(245, 222, 179),
                "title": RGBColor(139, 69, 19),
                "body": RGBColor(105, 105, 105),
                "title_font": "Georgia",
                "body_font": "Times New Roman",
            },
            "corporate": {
                "background": RGBColor(255, 255, 255),  # White background
                "title": RGBColor(0, 51, 102),  # Navy blue title
                "body": RGBColor(51, 51, 51),  # Dark gray body text
                "title_font": "Arial",
                "body_font": "Verdana",
            },
            "minimal": {
                "background": RGBColor(240, 240, 240),  # Light gray background
                "title": RGBColor(50, 50, 50),  # Dark gray title
                "body": RGBColor(80, 80, 80),  # Slightly lighter gray for body
                "title_font": "Helvetica",
                "body_font": "Sans-Serif",
            },
            "bold": {
                "background": RGBColor(0, 0, 0),  # Black background
                "title": RGBColor(255, 0, 0),  # Red title text
                "body": RGBColor(255, 255, 255),  # White body text
                "title_font": "Impact",
                "body_font": "Arial Black",
            },
        }
        self.theme = themes.get(name, themes["minimal"])


GENERATED_FOLDER = "gen/"
UPLOAD_FOLDER = "upload/"

app = Flask(__name__)
app.config["UPLOAD_FOLDER"] = UPLOAD_FOLDER
cors = CORS(app)

"""gTTS languages options:
    hi -> Hindi
    en -> English
    fr -> French
    
    Returns:
        _type_: audio file
"""

engine = pyttsx3.init()
engine.setProperty("rate", 170)
engine.setProperty("volume", 1)

MALE_VOICE = (
    "HKEY_LOCAL_MACHINE\SOFTWARE\Microsoft\Speech\Voices\Tokens\TTS_MS_EN-US_DAVID_11.0"
)
FEMALE_VOICE = (
    "HKEY_LOCAL_MACHINE\SOFTWARE\Microsoft\Speech\Voices\Tokens\TTS_MS_EN-US_ZIRA_11.0"
)


def concatenate_audio_wave(audio_clip_paths, output_path):
    data = []
    for clip in audio_clip_paths:
        w = wave.open(clip, "rb")
        data.append([w.getparams(), w.readframes(w.getnframes())])
        w.close()
    output = wave.open(output_path, "wb")
    output.setparams(data[0][0])
    for i in range(len(data)):
        output.writeframes(data[i][1])
    output.close()


@app.route("/generate-audio", methods=["POST"])
def generate_audio():
    if request.method == "POST":
        # req = request.get_json()
        # dialogues = req.get("dialogues")
        # random_hex = random.randint(1, 99999)
        # seq_number = 0
        # for dialogue in dialogues:
        #     speaker = list(dialogue.keys())[0]
        #     filename = f"audio/{random_hex}_{str(seq_number).zfill(5) }.mp3"

        #     if speaker == "katherine":
        #         engine.setProperty("voice", FEMALE_VOICE)
        #         engine.save_to_file(dialogue[speaker], filename)
        #         engine.runAndWait()
        #     else:
        #         engine.setProperty("voice", MALE_VOICE)
        #         engine.save_to_file(dialogue[speaker], filename)
        #         engine.runAndWait()

        #     seq_number += 1

        # filename = f"audio/merged_{random_hex}.mp3"
        # files = sorted(list(glob.glob(f"audio/{random_hex}_*")))

        # zip = zipfile.ZipFile(
        #     f"{GENERATED_FOLDER}/{random_hex}.zip", "w", zipfile.ZIP_DEFLATED
        # )

        # for file in files:
        #     zip.write(file)
        #     os.remove(file)

        # zip.close()

        # return_data = io.BytesIO()
        # with open(f"{GENERATED_FOLDER}/{random_hex}.zip", "rb") as fo:
        #     return_data.write(fo.read())
        # return_data.seek(0)

        # return send_file(
        #     return_data,
        #     mimetype="application/zip",
        # )

        return_data = io.BytesIO()
        with open(f"{GENERATED_FOLDER}/{"25579"}.zip", "rb") as fo:
            return_data.write(fo.read())
        return_data.seek(0)

        return send_file(
            return_data,
            mimetype="application/zip",
        )


def apply_background(slide, color):
    """Apply background color to a slide"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def create_ppt_from_dict(
    ppt_data: dict,
    image_mapping: dict,
    theme_name: str = "default",
    output_file: str = "presentation.pptx",
):
    prs = Presentation()
    theme = ThemeConfig(theme_name)

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Title Slide Fix
    title_slide_layout = prs.slide_layouts[0]  # Title slide layout
    title_slide = prs.slides.add_slide(title_slide_layout)
    apply_background(title_slide, theme.theme["background"])
    # Set title
    title = title_slide.shapes.title
    title.text = ppt_data["title"]
    title_para = title.text_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.name = theme.theme["title_font"]
    title_para.font.color.rgb = theme.theme["title"]
    title_para.alignment = PP_ALIGN.CENTER

    # Set subtitle (authors and institutions)
    subtitle = title_slide.placeholders[1]
    subtitle.text_frame.clear()  # Clear default placeholder text

    # Add authors as one paragraph
    authors_para = subtitle.text_frame.add_paragraph()
    authors_para.text = ", ".join(ppt_data["authors"])
    authors_para.font.size = Pt(18)
    authors_para.font.name = theme.theme["body_font"]
    authors_para.font.color.rgb = theme.theme["body"]
    authors_para.alignment = PP_ALIGN.CENTER

    # Add institution as a separate paragraph
    institution_para = subtitle.text_frame.add_paragraph()
    institution_para.text = "".join(ppt_data["institution"])
    institution_para.font.size = Pt(16)  # Slightly smaller font
    institution_para.font.name = theme.theme["body_font"]
    institution_para.font.color.rgb = theme.theme["body"]
    institution_para.alignment = PP_ALIGN.CENTER

    # Ensure the text fits within the shape
    subtitle.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    subtitle.text_frame.word_wrap = True

    # Add content slides
    for i in range(1, len(ppt_data["slides"])):
        slide_data = ppt_data["slides"][i]
        title_text = slide_data.get("title", "")

        # Detect Graphics/Graphs Slide
        is_graphics_slide = (
            "graphics" in title_text.lower() or "graphs slide" in title_text.lower()
        )

        # Use a blank layout for Graphics slides
        slide_layout = (
            prs.slide_layouts[6] if is_graphics_slide else prs.slide_layouts[1]
        )
        slide = prs.slides.add_slide(slide_layout)
        apply_background(slide, theme.theme["background"])

        if not is_graphics_slide:
            title = slide.shapes.title
            title.text = title_text
            title_para = title.text_frame.paragraphs[0]
            title_para.font.size = Pt(32)
            title_para.font.name = theme.theme["title_font"]
            title_para.font.color.rgb = theme.theme["title"]
        # Handling Graphics/Graphs Slide

        if is_graphics_slide and "images" in slide_data:
            image_filenames = slide_data["images"]
            image_paths = [
                image_mapping.get(fig.replace(".", "").replace(" ", ""))
                for fig in image_filenames
            ]
            image_paths = [
                img for img in image_paths if img and os.path.exists(img)
            ]  # Remove missing files

            num_images = len(image_paths)

            # Get theme colors
            caption_font = theme.theme["body_font"]
            caption_color = theme.theme["body"]

            # Define positioning based on number of images
            if num_images == 1:
                left, top, width, height = (
                    Inches(1.5),
                    Inches(1.5),
                    Inches(7),
                    Inches(5),
                )
                img_shape = slide.shapes.add_picture(
                    image_paths[0], left, top, width=width, height=height
                )
                caption_left = left + width / 2 - Inches(0.5)
                caption_top = top + height + Inches(0.2)

                # Add caption
                caption = slide.shapes.add_textbox(
                    caption_left, caption_top, Inches(1), Inches(0.5)
                )
                text_frame = caption.text_frame
                text_frame.text = image_filenames[0]
                para = text_frame.paragraphs[0]
                para.font.size = Pt(14)
                para.font.name = caption_font
                para.font.color.rgb = caption_color
                para.alignment = PP_ALIGN.CENTER

            elif num_images == 2:
                positions = [(Inches(1), Inches(2)), (Inches(5.5), Inches(2))]
                size = (Inches(4), Inches(3))

                for i, img_path in enumerate(image_paths[:2]):
                    img_left, img_top = positions[i]
                    img_shape = slide.shapes.add_picture(
                        img_path, img_left, img_top, *size
                    )

                    # Add caption
                    caption_left = img_left + size[0] / 2 - Inches(0.5)
                    caption_top = img_top + size[1] + Inches(0.2)
                    caption = slide.shapes.add_textbox(
                        caption_left, caption_top, Inches(1), Inches(0.5)
                    )
                    text_frame = caption.text_frame
                    text_frame.text = image_filenames[i]
                    para = text_frame.paragraphs[0]
                    para.font.size = Pt(14)
                    para.font.name = caption_font
                    para.font.color.rgb = caption_color
                    para.alignment = PP_ALIGN.CENTER

            elif num_images >= 3:
                positions = [
                    (Inches(1), Inches(1.5)),
                    (Inches(5), Inches(1.5)),
                    (Inches(3), Inches(4)),
                ]
                size = (Inches(3.5), Inches(2.5))

                for i, img_path in enumerate(image_paths[:3]):
                    img_left, img_top = positions[i]
                    img_shape = slide.shapes.add_picture(
                        img_path, img_left, img_top, *size
                    )

                    # Add caption
                    caption_left = img_left + size[0] / 2 - Inches(0.5)
                    caption_top = img_top + size[1] + Inches(0.2)
                    caption = slide.shapes.add_textbox(
                        caption_left, caption_top, Inches(1), Inches(0.5)
                    )
                    text_frame = caption.text_frame
                    text_frame.text = image_filenames[i]
                    para = text_frame.paragraphs[0]
                    para.font.size = Pt(14)
                    para.font.name = caption_font
                    para.font.color.rgb = caption_color
                    para.alignment = PP_ALIGN.CENTER

        else:
            bullet_points = slide_data.get("bullet_points", [])
            content_placeholder = slide.placeholders[1]
            text_frame = content_placeholder.text_frame
            text_frame.clear()
            if bullet_points:
                text_frame = content_placeholder.text_frame
                text_frame.clear()  # Remove default placeholder text
                text_frame.word_wrap = True  # Enable text wrapping
                text_frame.auto_size = (
                    MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
                )  # Enable auto size for content

                # Set default font size based on slide type
                is_references = "references" in slide_data.get("title", "").lower()
                DEFAULT_FONT_SIZE = 12 if is_references else 20

            for point in slide_data["bullet_points"]:
                paragraph = text_frame.add_paragraph()
                paragraph.text = point
                paragraph.font.size = Pt(DEFAULT_FONT_SIZE)
                paragraph.font.name = theme.theme["body_font"]
                paragraph.font.color.rgb = theme.theme["body"]

    # Save PowerPoint file
    prs.save(output_file)
    print(f"PowerPoint presentation saved as {output_file}")


@app.route("/hello")
def hello():
    return jsonify({"hello": "world"})


@app.route("/generate-ppt", methods=["POST"])
def generate_ppt():
    if request.method == "POST":
        # for file_name in request.files:
        #     print("HERE")
        #     file = request.files[file_name]

        #     print(file.filename)

        #     filename = secure_filename(file.filename)
        #     file.save(os.path.join(app.config["UPLOAD_FOLDER"], filename))

        #     url = "http://172.20.10.2:8000/generate-ppt"
        #     files = {
        #         "file": open(os.path.join(app.config["UPLOAD_FOLDER"], filename), "rb")
        #     }

        #     random_hex = random.randint(1, 999999)

        #     response = requests.post(url, files=files)
        #     structure = response.json()

        #     print(structure)

        #     create_ppt_from_dict(
        #         structure,
        #         image_mapping=[],
        #         theme_name="modern",
        #         output_file=f"{GENERATED_FOLDER}/{random_hex}.pptx",
        #     )

        #     return_data = io.BytesIO()
        #     with open(f"{GENERATED_FOLDER}/{random_hex}.pptx", "rb") as fo:
        #         return_data.write(fo.read())
        #     return_data.seek(0)

        #     return send_file(
        #         return_data,
        #         mimetype="application/msword",
        #     )
        return_data = io.BytesIO()
        with open(f"{GENERATED_FOLDER}/{"920087"}.pptx", "rb") as fo:
            return_data.write(fo.read())
        return_data.seek(0)

        return send_file(
            return_data,
            mimetype="application/msword",
        )


@app.route("/convert-pdf", methods=["POST"])
def convert_pdf():
    if request.method == "POST":
        for file_name in request.files:
            file = request.files[file_name]
            filename = secure_filename(file.filename)
            file.save(os.path.join(app.config["UPLOAD_FOLDER"], filename))

        return jsonify({"success": True})


if __name__ == "__main__":
    port = int(os.environ.get("PORT", 5000))
    app.run(host="0.0.0.0", port=port)
