from pptx import Presentation
from pptx.util import Inches


def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[0]  # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    title_shape.text = title
    subtitle_shape.text = subtitle


def add_bullet_slide(prs, title, bullet_points):
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    body_shape = slide.placeholders[1]
    title_shape.text = title
    tf = body_shape.text_frame
    tf.text = bullet_points[0]  # First bullet
    for point in bullet_points[1:]:
        p = tf.add_paragraph()
        p.text = point


def add_table_slide(prs, title, table_data):
    slide_layout = prs.slide_layouts[5]  # Title Only
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    title_shape.text = title
    rows, cols = len(table_data), len(table_data[0])
    table = slide.shapes.add_table(
        rows, cols, Inches(1), Inches(1.5), Inches(8), Inches(3)
    ).table
    for r, row in enumerate(table_data):
        for c, cell in enumerate(row):
            table.cell(r, c).text = cell


def create_presentation():
    prs = Presentation()

    # Slide 1: Title Slide
    add_title_slide(
        prs,
        "Transformer: A Novel Neural Network Architecture for Machine Translation",
        "Authors: Vaswani et al.\nDate: 2017",
    )

    # Bullet Point Slides
    slides = [
        (
            "Introduction",
            [
                "Machine translation is a challenging task that requires understanding input sentences.",
                "Traditional sequence-to-sequence models have limitations.",
                "The Transformer addresses these limitations.",
            ],
        ),
        (
            "Architecture Overview",
            [
                "The Transformer consists of an encoder and a decoder.",
                "The encoder processes tokens into vectors.",
                "The decoder generates output tokens.",
            ],
        ),
        (
            "Encoder",
            [
                "Consists of stacked identical layers.",
                "Each layer includes:",
                "- Self-attention mechanism",
                "- Position-wise feed-forward network",
            ],
        ),
        (
            "Self-Attention Mechanism",
            [
                "Allows the model to attend to all positions in the input sequence.",
                "Key components:",
                "- Query (Q)",
                "- Key (K)",
                "- Value (V)",
            ],
        ),
        (
            "Training",
            [
                "Uses Adam optimizer with a scheduled learning rate.",
                "Trained on WMT 2014 English-German dataset.",
            ],
        ),
        (
            "Results",
            [
                "Achieves state-of-the-art results in translation tasks.",
                "Outperforms previous models.",
            ],
        ),
        (
            "Future Work",
            [
                "Extending the Transformer to other tasks.",
                "Investigating local attention mechanisms.",
            ],
        ),
    ]

    for title, content in slides:
        add_bullet_slide(prs, title, content)

    # Table Slides
    tables = [
        (
            "Architecture Overview",
            [
                ["Component", "Description"],
                ["Encoder", "Processes input sequence"],
                ["Decoder", "Generates output sequence"],
                ["Self-Attention", "Captures dependencies across sequence"],
            ],
        ),
        (
            "Training Details",
            [
                ["Component", "Description"],
                ["Adam Optimizer", "Used for training"],
                ["Learning Rate Schedule", "Linear increase, then decrease"],
            ],
        ),
        (
            "Results",
            [
                ["Task", "BLEU Score"],
                ["English-German", "28.4"],
                ["English-French", "41.0"],
            ],
        ),
    ]

    for title, table_data in tables:
        add_table_slide(prs, title, table_data)

    prs.save("transformer_presentation.pptx")


create_presentation()
